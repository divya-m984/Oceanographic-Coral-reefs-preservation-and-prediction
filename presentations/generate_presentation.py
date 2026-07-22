#!/usr/bin/env python3
"""
Generate Oceanographic Presentation 2 — Data Collection and Model Training.

Output: presentations/Oceanographic_Presentation_2_Data_Collection_and_Model_Training.pptx

All metrics sourced from:
  - models/evaluation_health.json
  - models/evaluation_restoration.json
  - params.yaml
  - data/raw/observations.csv (column count verified)

Run from project root:
  .venv/bin/python presentations/generate_presentation.py
"""

import io
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).parent.parent
OUT_DIR = ROOT / "presentations"
OUT_FILE = OUT_DIR / "Oceanographic_Presentation_2_Data_Collection_and_Model_Training.pptx"

# ─────────────────────────────────────────────────────────────────────────────
# Confirmed repository data (do not modify)
# ─────────────────────────────────────────────────────────────────────────────
TOTAL_SAMPLES   = 15_000
CSV_COLUMNS     = 21        # confirmed by: wc -l + python csv header count
TRAIN_ROWS      = 12_000    # 80% of 15,000 (test_size=0.20 in params.yaml)
TEST_ROWS       = 3_000     # 20% of 15,000
CV_FOLDS        = 5         # models.cv_folds in params.yaml
RANDOM_SEED     = 42        # base.random_seed in params.yaml
NOISE_SCALE     = 0.15      # data.noise_scale in params.yaml

# Confirmed from models/evaluation_health.json
HEALTH_METRICS = {
    "Logistic\nRegression": {
        "cv_macro_f1":  0.7612,
        "test_macro_f1": 0.7871,
        "balanced_acc": 0.8012,
        "accuracy":     0.8193,
        "selected": True,
    },
    "Random\nForest": {
        "cv_macro_f1":  0.7605,
        "test_macro_f1": 0.7778,
        "balanced_acc": 0.7873,
        "accuracy":     0.8143,
        "selected": False,
    },
    "XGBoost": {
        "cv_macro_f1":  0.7580,
        "test_macro_f1": 0.7740,
        "balanced_acc": 0.7804,
        "accuracy":     0.8157,
        "selected": False,
    },
}

# Confirmed from models/evaluation_restoration.json
RESTORATION_METRICS = {
    "Logistic\nRegression": {
        "cv_macro_f1":  0.7808,
        "test_macro_f1": 0.7863,
        "balanced_acc": 0.8163,
        "accuracy":     0.8177,
        "selected": False,
    },
    "Random\nForest": {
        "cv_macro_f1":  0.7839,
        "test_macro_f1": 0.7918,
        "balanced_acc": 0.8108,
        "accuracy":     0.8193,
        "selected": False,
    },
    "XGBoost": {
        "cv_macro_f1":  0.7913,
        "test_macro_f1": 0.8029,
        "balanced_acc": 0.8121,
        "accuracy":     0.8277,
        "selected": True,
    },
}

# ─────────────────────────────────────────────────────────────────────────────
# Colour palette — deep navy oceanographic theme
# ─────────────────────────────────────────────────────────────────────────────
C_NAVY       = RGBColor(0x08, 0x12, 0x24)
C_NAVY2      = RGBColor(0x0D, 0x1F, 0x3C)
C_NAVY3      = RGBColor(0x12, 0x2A, 0x4A)
C_TEAL       = RGBColor(0x1B, 0x8A, 0x8F)
C_TEAL_DARK  = RGBColor(0x0E, 0x5A, 0x5E)
C_CYAN       = RGBColor(0x00, 0xB4, 0xD8)
C_CYAN_LIGHT = RGBColor(0x90, 0xE0, 0xEF)
C_CORAL      = RGBColor(0xFF, 0x6B, 0x35)
C_AMBER      = RGBColor(0xF4, 0xA2, 0x61)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT       = RGBColor(0xE8, 0xF4, 0xF8)
C_MUTED      = RGBColor(0x8A, 0xB4, 0xC8)
C_HIGHLIGHT  = RGBColor(0x14, 0x3D, 0x60)
C_GREEN      = RGBColor(0x2E, 0xC4, 0xB6)

# Matplotlib equivalents
M_BG     = "#081224"
M_NAVY2  = "#0D1F3C"
M_NAVY3  = "#122A4A"
M_TEAL   = "#1B8A8F"
M_CYAN   = "#00B4D8"
M_CORAL  = "#FF6B35"
M_AMBER  = "#F4A261"
M_WHITE  = "#E8F4F8"
M_MUTED  = "#8AB4C8"
M_GREEN  = "#2EC4B6"
M_GOLD   = "#FFD700"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FOOTER_TEXT = "Oceanographic Reef Prediction Framework"


# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color=C_NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill=C_NAVY2, line=None, line_w=Pt(0)):
    """Add a plain rectangle shape."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_tb(slide, left, top, width, height, text="", size=Pt(14), bold=False,
           color=C_TEXT, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a textbox with a single paragraph."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_tb_multi(slide, left, top, width, height, lines, size=Pt(13),
                 color=C_TEXT, align=PP_ALIGN.LEFT, line_spacing=None):
    """Add a textbox with multiple paragraphs (list of (text, bold, size, color))."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    from lxml import etree
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, fsize, fcolor = item, False, size, color
        elif len(item) == 2:
            text, bold = item
            fsize, fcolor = size, color
        elif len(item) == 3:
            text, bold, fsize = item
            fcolor = color
        else:
            text, bold, fsize, fcolor = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = fsize
        run.font.bold = bold
        run.font.color.rgb = fcolor
    return txb


def add_footer(slide, slide_num: int):
    """Add slide number (right) and footer text (left)."""
    # Footer bar line (thin teal)
    add_shape(slide,
              left=Inches(0), top=Inches(7.18),
              width=SLIDE_W, height=Inches(0.04),
              fill=C_TEAL)
    # Footer text left
    add_tb(slide,
           left=Inches(0.3), top=Inches(7.22),
           width=Inches(9), height=Inches(0.25),
           text=FOOTER_TEXT,
           size=Pt(9), color=C_MUTED, align=PP_ALIGN.LEFT)
    # Slide number right
    add_tb(slide,
           left=Inches(12.3), top=Inches(7.22),
           width=Inches(0.8), height=Inches(0.25),
           text=str(slide_num),
           size=Pt(9), color=C_MUTED, align=PP_ALIGN.RIGHT)


def add_accent_line(slide, color=C_TEAL):
    """Thin accent line just below the header area."""
    add_shape(slide,
              left=Inches(0), top=Inches(1.12),
              width=SLIDE_W, height=Inches(0.03),
              fill=color)


def add_slide_title(slide, title: str, subtitle: str = ""):
    """Standard slide title block."""
    # Title background
    add_shape(slide,
              left=Inches(0), top=Inches(0),
              width=SLIDE_W, height=Inches(1.12),
              fill=C_NAVY2)
    # Title text
    add_tb(slide,
           left=Inches(0.35), top=Inches(0.15),
           width=Inches(12.5), height=Inches(0.65),
           text=title,
           size=Pt(24), bold=True, color=C_CYAN, align=PP_ALIGN.LEFT)
    if subtitle:
        add_tb(slide,
               left=Inches(0.35), top=Inches(0.72),
               width=Inches(12.5), height=Inches(0.35),
               text=subtitle,
               size=Pt(12), color=C_MUTED, align=PP_ALIGN.LEFT)
    add_accent_line(slide)


def png_bytes(fig) -> bytes:
    """Render matplotlib figure to PNG bytes."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor(), edgecolor="none")
    buf.seek(0)
    data = buf.read()
    buf.close()
    plt.close(fig)
    return data


def add_image_bytes(slide, img_bytes: bytes, left, top, width, height):
    """Add PNG bytes as an image on the slide."""
    buf = io.BytesIO(img_bytes)
    slide.shapes.add_picture(buf, left, top, width, height)


def add_notes(slide, text: str):
    """Set speaker notes text."""
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ─────────────────────────────────────────────────────────────────────────────
# Chart builders
# ─────────────────────────────────────────────────────────────────────────────

def make_model_comparison_chart(metrics: dict, task_label: str) -> bytes:
    """
    Grouped bar chart comparing three algorithms across four metrics.
    Selected model highlighted with gold outline and star marker.
    """
    models    = list(metrics.keys())
    metric_labels = ["CV Macro-F1", "Test Macro-F1", "Balanced Acc.", "Accuracy"]
    metric_keys   = ["cv_macro_f1", "test_macro_f1", "balanced_acc", "accuracy"]

    # Base colours: teal for unselected, cyan+gold for selected
    bar_colors = []
    for m in models:
        if metrics[m]["selected"]:
            bar_colors.append([M_CYAN, M_GREEN, M_TEAL, "#5BC0EB"])
        else:
            bar_colors.append(["#1A4A6E", "#1E5870", "#173B52", "#143248"])

    n_models  = len(models)
    n_metrics = len(metric_keys)
    x = np.arange(n_metrics)
    bar_w = 0.24
    offsets = np.array([-1, 0, 1]) * bar_w

    fig, ax = plt.subplots(figsize=(9.5, 4.4))
    fig.patch.set_facecolor(M_BG)
    ax.set_facecolor(M_NAVY3)

    rects_all = []
    for i, (model, offset) in enumerate(zip(models, offsets)):
        vals = [metrics[model][k] for k in metric_keys]
        rects = ax.bar(x + offset, vals, bar_w - 0.02,
                       color=bar_colors[i], zorder=3,
                       edgecolor="none")
        rects_all.append((model, rects, metrics[model]["selected"]))

    # Highlight selected model with gold border on bars
    for model, rects, selected in rects_all:
        if selected:
            for rect in rects:
                rect.set_edgecolor(M_GOLD)
                rect.set_linewidth(1.6)

    # Value labels on bars
    for model, rects, selected in rects_all:
        for rect in rects:
            h = rect.get_height()
            ax.text(rect.get_x() + rect.get_width() / 2.0,
                    h + 0.005,
                    f"{h:.4f}",
                    ha="center", va="bottom",
                    fontsize=6.5,
                    color=M_GOLD if selected else M_MUTED,
                    fontweight="bold" if selected else "normal")

    # Y axis
    ax.set_ylim(0.72, 0.88)
    ax.set_yticks(np.arange(0.72, 0.89, 0.02))
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{v:.2f}"))
    ax.tick_params(axis="y", colors=M_MUTED, labelsize=8)
    ax.set_yticklabels([f"{v:.2f}" for v in np.arange(0.72, 0.89, 0.02)],
                       color=M_MUTED, fontsize=8)

    # X axis
    ax.set_xticks(x)
    ax.set_xticklabels(metric_labels, color=M_WHITE, fontsize=10, fontweight="medium")
    ax.tick_params(axis="x", length=0)

    # Grid
    ax.grid(axis="y", color="#1E3A5F", linewidth=0.8, zorder=0)
    ax.set_axisbelow(True)

    # Spines
    for spine in ax.spines.values():
        spine.set_visible(False)

    # Legend
    legend_handles = []
    for i, (model, rects, selected) in enumerate(rects_all):
        display = model.replace("\n", " ")
        patch = mpatches.Patch(
            color=bar_colors[i][0],
            label=f"{'* ' if selected else ''}{display}{'  (selected)' if selected else ''}",
            linewidth=2.0 if selected else 0,
            edgecolor=M_GOLD if selected else "none"
        )
        legend_handles.append(patch)

    leg = ax.legend(handles=legend_handles,
                    loc="lower right",
                    frameon=True,
                    framealpha=0.85,
                    facecolor=M_NAVY2,
                    edgecolor=M_TEAL,
                    fontsize=9,
                    labelcolor=M_WHITE)

    ax.set_ylabel("Score", color=M_MUTED, fontsize=9)
    ax.set_title(task_label, color=M_WHITE, fontsize=11, fontweight="bold", pad=10)

    # Selected model annotation
    for model, data in metrics.items():
        if data["selected"]:
            sel_label = model.replace("\n", " ")
            ax.annotate(f"Selected: {sel_label}",
                        xy=(0.02, 0.97), xycoords="axes fraction",
                        fontsize=8.5, color=M_GOLD, fontweight="bold",
                        va="top")

    fig.tight_layout(pad=0.8)
    return png_bytes(fig)


def make_preprocessing_diagram() -> bytes:
    """Flow diagram for the preprocessing pipeline (Slide 7)."""
    fig, ax = plt.subplots(figsize=(11, 4.0))
    fig.patch.set_facecolor(M_BG)
    ax.set_facecolor(M_BG)
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 4)
    ax.axis("off")

    # Top row: main pipeline flow
    top_boxes = [
        (0.5,  2.8, "Raw\nObservations", M_NAVY3, M_CYAN),
        (2.3,  2.8, "Separate\nInputs & Targets", M_NAVY3, M_TEAL),
        (4.2,  2.8, "Stratified\n80:20 Split", M_NAVY3, M_TEAL),
        (6.1,  2.8, "Fit Preprocessor\n(train only)", M_NAVY3, M_CORAL),
        (8.0,  2.8, "Transform\nTrain & Test", M_NAVY3, M_TEAL),
        (9.9,  2.8, "Save Artifacts\n(.joblib)", M_NAVY3, M_GREEN),
    ]

    bw, bh = 1.55, 0.85
    for (bx, by, label, bg, border) in top_boxes:
        box = FancyBboxPatch((bx, by), bw, bh,
                              boxstyle="round,pad=0.06",
                              facecolor=bg, edgecolor=border,
                              linewidth=1.8, zorder=3)
        ax.add_patch(box)
        ax.text(bx + bw / 2, by + bh / 2, label,
                ha="center", va="center",
                fontsize=8.5, color=M_WHITE,
                fontweight="medium", zorder=4,
                multialignment="center")

    # Arrows between top row boxes
    for i in range(len(top_boxes) - 1):
        x0 = top_boxes[i][0] + bw
        x1 = top_boxes[i + 1][0]
        y  = top_boxes[i][1] + bh / 2
        ax.annotate("", xy=(x1, y), xytext=(x0, y),
                    arrowprops=dict(arrowstyle="-|>",
                                   color=M_MUTED,
                                   lw=1.5))

    # Bottom row: transformation details
    bottom_boxes = [
        (1.2, 1.1, "Numeric Features (15)\nStandardScaler\n+ Median Imputer", M_NAVY3, M_CYAN),
        (4.2, 1.1, "Categorical Feature (1)\nOneHotEncoder\n(region, 4 values)", M_NAVY3, M_TEAL),
        (7.2, 1.1, "6 Derived Features\nEngineered from\nraw sensor readings", M_NAVY3, M_AMBER),
        (9.9, 1.1, "Leakage Guard\nPreprocessor fitted\non training split only", M_NAVY3, M_CORAL),
    ]
    bw2, bh2 = 2.6, 1.0

    for (bx, by, label, bg, border) in bottom_boxes:
        box = FancyBboxPatch((bx, by), bw2, bh2,
                              boxstyle="round,pad=0.06",
                              facecolor=bg, edgecolor=border,
                              linewidth=1.5, zorder=3)
        ax.add_patch(box)
        ax.text(bx + bw2 / 2, by + bh2 / 2, label,
                ha="center", va="center",
                fontsize=8, color=M_WHITE,
                fontweight="normal", zorder=4,
                multialignment="center")

    # Connecting dashed line from top to bottom
    ax.annotate("", xy=(2.5, 1.1 + bh2), xytext=(2.5, top_boxes[0][1]),
                arrowprops=dict(arrowstyle="-|>", color="#4A6A8A", lw=1.2,
                                linestyle="dashed"))
    ax.annotate("", xy=(5.5, 1.1 + bh2), xytext=(5.5, top_boxes[0][1]),
                arrowprops=dict(arrowstyle="-|>", color="#4A6A8A", lw=1.2,
                                linestyle="dashed"))

    fig.tight_layout(pad=0.3)
    return png_bytes(fig)


def make_pipeline_diagram() -> bytes:
    """DVC pipeline stages diagram for Slide 11."""
    fig, ax = plt.subplots(figsize=(11, 2.5))
    fig.patch.set_facecolor(M_BG)
    ax.set_facecolor(M_BG)
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 2.5)
    ax.axis("off")

    stages = [
        ("Generate",          M_TEAL),
        ("Validate\n(Pandera)", M_TEAL),
        ("Preprocess",        M_TEAL),
        ("Train\n(LR/RF/XGB)", M_CYAN),
        ("Evaluate",          M_TEAL),
        ("Register\nCandidate", M_AMBER),
    ]
    n = len(stages)
    gap = 14.0 / n
    bw, bh = gap * 0.72, 1.1

    for i, (label, color) in enumerate(stages):
        bx = i * gap + (gap - bw) / 2
        by = 0.65
        box = FancyBboxPatch((bx, by), bw, bh,
                              boxstyle="round,pad=0.08",
                              facecolor=M_NAVY3, edgecolor=color,
                              linewidth=2.2, zorder=3)
        ax.add_patch(box)
        ax.text(bx + bw / 2, by + bh / 2, label,
                ha="center", va="center",
                fontsize=9, color=M_WHITE,
                fontweight="medium", zorder=4,
                multialignment="center")

        # Arrow to next
        if i < n - 1:
            x0 = bx + bw
            x1 = (i + 1) * gap + (gap - bw) / 2
            y  = by + bh / 2
            ax.annotate("", xy=(x1, y), xytext=(x0, y),
                        arrowprops=dict(arrowstyle="-|>",
                                       color=M_MUTED, lw=1.8))

    # Note: candidate ≠ champion
    ax.text(7.0, 0.22,
            "Register Candidate does NOT set the champion alias — explicit approval required",
            ha="center", va="bottom", fontsize=8.5,
            color=M_AMBER, style="italic")

    fig.tight_layout(pad=0.3)
    return png_bytes(fig)


def make_dataset_summary_chart() -> bytes:
    """Compact dataset summary visualization for Slide 5."""
    fig, axes = plt.subplots(1, 2, figsize=(9.5, 3.2))
    fig.patch.set_facecolor(M_BG)

    # Health class distribution (approximate from evaluation_health.json test support)
    # test support: healthy=318 (10.6%), stressed=1726 (57.5%), bleached=283 (9.4%), severely_degraded=673 (22.4%)
    # Scaling back to 15,000: roughly proportional
    health_labels  = ["healthy", "stressed", "bleached", "severely\ndegraded"]
    # From test set proportions: 318/3000*15000 ≈ 1590, 1726/3000*15000 ≈ 8630, 283/3000*15000 ≈ 1415, 673/3000*15000 ≈ 3365
    health_counts  = [1590, 8630, 1415, 3365]
    health_colors  = [M_GREEN, M_AMBER, M_CORAL, "#8B0000"]

    ax1 = axes[0]
    ax1.set_facecolor(M_NAVY3)
    bars1 = ax1.bar(health_labels, health_counts,
                    color=health_colors, edgecolor=M_NAVY3, linewidth=0.5, zorder=3)
    ax1.set_title("Reef Health Classes", color=M_WHITE, fontsize=10, pad=6)
    ax1.set_ylabel("Approx. count", color=M_MUTED, fontsize=8)
    ax1.tick_params(axis="x", colors=M_WHITE, labelsize=7.5)
    ax1.tick_params(axis="y", colors=M_MUTED, labelsize=7.5)
    ax1.grid(axis="y", color="#1E3A5F", linewidth=0.6, zorder=0)
    ax1.set_axisbelow(True)
    for spine in ax1.spines.values():
        spine.set_visible(False)
    for bar in bars1:
        h = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width() / 2, h + 80,
                 f"{h:,}", ha="center", va="bottom",
                 fontsize=7, color=M_WHITE)

    # Restoration class distribution
    # test support: suitable=907, moderately_suitable=1798, unsuitable=295
    # Scale: 907/3000*15000≈4535, 1798/3000*15000≈8990, 295/3000*15000≈1475
    rest_labels = ["suitable", "moderately\nsuitable", "unsuitable"]
    rest_counts = [4535, 8990, 1475]
    rest_colors = [M_GREEN, M_TEAL, M_CORAL]

    ax2 = axes[1]
    ax2.set_facecolor(M_NAVY3)
    bars2 = ax2.bar(rest_labels, rest_counts,
                    color=rest_colors, edgecolor=M_NAVY3, linewidth=0.5, zorder=3)
    ax2.set_title("Restoration Suitability Classes", color=M_WHITE, fontsize=10, pad=6)
    ax2.set_ylabel("Approx. count", color=M_MUTED, fontsize=8)
    ax2.tick_params(axis="x", colors=M_WHITE, labelsize=8)
    ax2.tick_params(axis="y", colors=M_MUTED, labelsize=7.5)
    ax2.grid(axis="y", color="#1E3A5F", linewidth=0.6, zorder=0)
    ax2.set_axisbelow(True)
    for spine in ax2.spines.values():
        spine.set_visible(False)
    for bar in bars2:
        h = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width() / 2, h + 80,
                 f"{h:,}", ha="center", va="bottom",
                 fontsize=7.5, color=M_WHITE)

    fig.suptitle("Class Distribution  (15,000 synthetic observations — approximate from test-split proportions)",
                 color=M_MUTED, fontsize=7.5, y=0.02)
    fig.tight_layout(pad=0.8)
    return png_bytes(fig)


def make_validation_flow() -> bytes:
    """Validation flow diagram for Slide 6."""
    fig, ax = plt.subplots(figsize=(11, 3.6))
    fig.patch.set_facecolor(M_BG)
    ax.set_facecolor(M_BG)
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 4)
    ax.axis("off")

    # Input box
    input_box = FancyBboxPatch((0.2, 1.55), 1.8, 0.9,
                                boxstyle="round,pad=0.06",
                                facecolor=M_NAVY3, edgecolor=M_CYAN,
                                linewidth=2.0, zorder=3)
    ax.add_patch(input_box)
    ax.text(1.1, 2.0, "Raw CSV\nInput", ha="center", va="center",
            fontsize=9, color=M_WHITE, fontweight="medium", multialignment="center")

    # Checks in a row
    checks = [
        ("Column\npresence", M_TEAL),
        ("Data types\n& coercion", M_TEAL),
        ("Sensor range\nbounds", M_TEAL),
        ("Missing\nvalue check", M_TEAL),
        ("Duplicate\nrow check", M_TEAL),
        ("Coordinate\nbounds", M_TEAL),
        ("Target label\nvalidity", M_TEAL),
    ]
    start_x = 2.5
    cw, ch = 1.45, 0.85
    gap_c = 1.65

    for i, (label, color) in enumerate(checks):
        cx = start_x + i * gap_c
        cy = 1.575
        box = FancyBboxPatch((cx, cy), cw, ch,
                              boxstyle="round,pad=0.05",
                              facecolor=M_NAVY3, edgecolor=color,
                              linewidth=1.5, zorder=3)
        ax.add_patch(box)
        ax.text(cx + cw / 2, cy + ch / 2, label,
                ha="center", va="center",
                fontsize=7.8, color=M_WHITE,
                multialignment="center")

    # Arrow from input to first check
    ax.annotate("", xy=(start_x, 2.0), xytext=(2.0, 2.0),
                arrowprops=dict(arrowstyle="-|>", color=M_MUTED, lw=1.5))

    # Arrows between checks
    for i in range(len(checks) - 1):
        x0 = start_x + i * gap_c + cw
        x1 = start_x + (i + 1) * gap_c
        ax.annotate("", xy=(x1, 2.0), xytext=(x0, 2.0),
                    arrowprops=dict(arrowstyle="-|>", color=M_MUTED, lw=1.2))

    # Outcome boxes
    x_last = start_x + (len(checks) - 1) * gap_c + cw
    # Pass path
    pass_box = FancyBboxPatch((x_last + 0.25, 2.35), 1.4, 0.65,
                               boxstyle="round,pad=0.05",
                               facecolor="#0A3A1A", edgecolor=M_GREEN,
                               linewidth=2.0, zorder=3)
    ax.add_patch(pass_box)
    ax.text(x_last + 0.95, 2.675, "PASS\nNext stage", ha="center",
            va="center", fontsize=8, color=M_GREEN, fontweight="bold",
            multialignment="center")

    # Fail path
    fail_box = FancyBboxPatch((x_last + 0.25, 1.05), 1.4, 0.65,
                               boxstyle="round,pad=0.05",
                               facecolor="#3A0A0A", edgecolor=M_CORAL,
                               linewidth=2.0, zorder=3)
    ax.add_patch(fail_box)
    ax.text(x_last + 0.95, 1.375, "FAIL\nPipeline stops", ha="center",
            va="center", fontsize=8, color=M_CORAL, fontweight="bold",
            multialignment="center")

    # Arrow to pass/fail
    ax.annotate("", xy=(x_last + 0.25, 2.675), xytext=(x_last, 2.675),
                arrowprops=dict(arrowstyle="-|>", color=M_GREEN, lw=1.5))
    ax.annotate("", xy=(x_last + 0.25, 1.375), xytext=(x_last, 1.375),
                arrowprops=dict(arrowstyle="-|>", color=M_CORAL, lw=1.5))
    # Vertical connector
    ax.plot([x_last, x_last], [1.375, 2.675], color=M_MUTED, lw=1.2)

    # Bottom note
    ax.text(7.0, 0.3,
            "Pandera v0.20+ DataFrameModel schema — lazy validation collects all errors before reporting",
            ha="center", va="bottom", fontsize=8, color=M_MUTED, style="italic")

    fig.tight_layout(pad=0.3)
    return png_bytes(fig)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)

    # Deep blue gradient panel (bottom half)
    add_shape(slide,
              left=Inches(0), top=Inches(3.9),
              width=SLIDE_W, height=Inches(3.6),
              fill=C_NAVY2)

    # Decorative teal band
    add_shape(slide,
              left=Inches(0), top=Inches(3.88),
              width=SLIDE_W, height=Inches(0.05),
              fill=C_TEAL)

    # Cyan accent vertical bar
    add_shape(slide,
              left=Inches(0), top=Inches(0),
              width=Inches(0.12), height=SLIDE_H,
              fill=C_TEAL)

    # Project title (full exact title, no abbreviation)
    title_text = (
        "Oceanographic: A Machine Learning-Driven Sonar Framework\n"
        "for Real-Time Coral Reef Habitat Prediction\n"
        "and Marine Ecosystem Monitoring"
    )
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(1.05), Inches(12.6), Inches(2.0))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = C_CYAN

    # Subtitle
    add_tb(slide,
           left=Inches(0.4), top=Inches(3.25),
           width=Inches(11), height=Inches(0.55),
           text="MLOps Presentation 2 — Data Collection and Model Training",
           size=Pt(16), color=C_CYAN_LIGHT, align=PP_ALIGN.LEFT)

    # Decorative rule
    add_shape(slide,
              left=Inches(0.4), top=Inches(3.88 + 0.18),
              width=Inches(11), height=Inches(0.025),
              fill=C_TEAL)

    # Bottom detail lines
    detail_lines = [
        "Decision-support prototype  |  Synthetic data  |  Indian reef systems",
        "Four regions: Lakshadweep  |  Gulf of Mannar  |  Gulf of Kutch  |  Andaman and Nicobar Islands",
    ]
    for i, line in enumerate(detail_lines):
        add_tb(slide,
               left=Inches(0.4), top=Inches(4.22 + i * 0.45),
               width=Inches(12.5), height=Inches(0.42),
               text=line,
               size=Pt(12), color=C_MUTED, align=PP_ALIGN.LEFT)

    # Disclaimer
    add_tb(slide,
           left=Inches(0.4), top=Inches(6.35),
           width=Inches(12.5), height=Inches(0.42),
           text="All observations in this prototype are synthetic. Models do not represent real-world conservation accuracy.",
           size=Pt(9.5), italic=True, color=RGBColor(0x7A, 0xA4, 0xB8), align=PP_ALIGN.LEFT)

    # Slide number only (no full footer on title)
    add_shape(slide,
              left=Inches(0), top=Inches(7.18),
              width=SLIDE_W, height=Inches(0.04),
              fill=C_TEAL)
    add_tb(slide,
           left=Inches(12.3), top=Inches(7.22),
           width=Inches(0.8), height=Inches(0.25),
           text="1",
           size=Pt(9), color=C_MUTED, align=PP_ALIGN.RIGHT)

    add_notes(slide,
        "Welcome. This presentation covers the data collection strategy and model training "
        "pipeline for the Oceanographic reef prediction framework. This is a machine learning "
        "prototype designed to classify reef health and assess restoration suitability from "
        "simulated marine sensor observations. I'll walk through how the data was generated, "
        "how it was processed and validated, how the models were trained, and what results "
        "were achieved on a synthetic benchmark dataset. Everything you see here is backed by "
        "actual code and evaluation files in the project repository."
    )
    return slide


def slide_02_problem(prs):
    """Problem statement and ML objectives."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Problem and ML Objectives",
                    "Why machine learning for coral reef monitoring?")
    add_footer(slide, 2)

    # Left column — problem
    add_shape(slide,
              left=Inches(0.3), top=Inches(1.3),
              width=Inches(5.9), height=Inches(5.6),
              fill=C_NAVY3, line=C_TEAL, line_w=Pt(0.8))

    add_tb(slide,
           left=Inches(0.5), top=Inches(1.38),
           width=Inches(5.5), height=Inches(0.38),
           text="The Monitoring Challenge",
           size=Pt(13), bold=True, color=C_CYAN)

    problem_lines = [
        "Coral reefs cover less than 1% of the ocean floor but support",
        "approximately 25% of all marine species. Monitoring reef health",
        "across large remote areas requires combining multiple sensor types:",
        "",
        "  Sonar — measures structural complexity (rugosity, backscatter,",
        "  hard substrate), not bleaching or coral cover directly.",
        "",
        "  Environmental sensors — capture temperature, pH, salinity,",
        "  dissolved oxygen, turbidity, and current speed.",
        "",
        "  Biological surveys — observe coral cover, bleaching extent,",
        "  and disease prevalence.",
        "",
        "Manual expert assessment at scale is costly and slow. A decision-",
        "support model can flag priority sites for human follow-up.",
    ]
    y = 1.82
    for line in problem_lines:
        add_tb(slide,
               left=Inches(0.48), top=Inches(y),
               width=Inches(5.55), height=Inches(0.28),
               text=line, size=Pt(10.5), color=C_TEXT)
        y += 0.26

    # Right column — ML tasks
    add_shape(slide,
              left=Inches(6.5), top=Inches(1.3),
              width=Inches(6.55), height=Inches(2.55),
              fill=C_NAVY3, line=C_CYAN, line_w=Pt(0.8))

    add_tb(slide,
           left=Inches(6.7), top=Inches(1.38),
           width=Inches(6.2), height=Inches(0.38),
           text="Task 1 — Reef Health Classification",
           size=Pt(13), bold=True, color=C_CYAN)

    health_lines = [
        "Classify each observation into one of four ordinal states:",
        "  healthy  |  stressed  |  bleached  |  severely degraded",
        "",
        "Input: 16 sensor and survey features per observation.",
        "Output: predicted reef health class + class probabilities.",
    ]
    y = 1.83
    for line in health_lines:
        add_tb(slide,
               left=Inches(6.68), top=Inches(y),
               width=Inches(6.2), height=Inches(0.28),
               text=line, size=Pt(10.5), color=C_TEXT)
        y += 0.26

    add_shape(slide,
              left=Inches(6.5), top=Inches(4.05),
              width=Inches(6.55), height=Inches(2.35),
              fill=C_NAVY3, line=C_TEAL, line_w=Pt(0.8))

    add_tb(slide,
           left=Inches(6.7), top=Inches(4.13),
           width=Inches(6.2), height=Inches(0.38),
           text="Task 2 — Restoration Suitability Classification",
           size=Pt(13), bold=True, color=C_TEAL)

    rest_lines = [
        "Classify each site into one of three restoration categories:",
        "  suitable  |  moderately suitable  |  unsuitable",
        "",
        "Informed by substrate quality, temperature stability,",
        "turbidity, depth, current, and existing coral proximity.",
    ]
    y = 4.57
    for line in rest_lines:
        add_tb(slide,
               left=Inches(6.68), top=Inches(y),
               width=Inches(6.2), height=Inches(0.28),
               text=line, size=Pt(10.5), color=C_TEXT)
        y += 0.26

    # Prototype note
    add_shape(slide,
              left=Inches(6.5), top=Inches(6.52),
              width=Inches(6.55), height=Inches(0.34),
              fill=C_HIGHLIGHT)
    add_tb(slide,
           left=Inches(6.6), top=Inches(6.54),
           width=Inches(6.35), height=Inches(0.3),
           text="This is a decision-support prototype — outputs guide human experts, not automated conservation policy.",
           size=Pt(9.5), italic=True, color=C_AMBER)

    add_notes(slide,
        "Coral reefs are under severe stress from ocean warming, acidification, and human activity. "
        "The challenge is that monitoring them at scale is hard — you need to combine sonar readings "
        "for structure, physical sensor data for water conditions, and direct biological surveys. "
        "This project frames two supervised classification tasks. The first predicts whether a reef "
        "is healthy, stressed, bleached, or severely degraded. The second predicts whether a site "
        "is suitable for active restoration intervention. Both tasks use the same sensor features "
        "as input. It is very important to note that this is a prototype — the model outputs are "
        "intended to assist expert ecologists, not replace their judgment."
    )
    return slide


def slide_03_data_collection(prs):
    """Data collection strategy."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Data Collection Strategy",
                    "Current prototype vs. proposed future deployment")
    add_footer(slide, 3)

    # Left: current (synthetic)
    add_shape(slide,
              left=Inches(0.3), top=Inches(1.3),
              width=Inches(5.95), height=Inches(5.6),
              fill=C_NAVY3, line=C_AMBER, line_w=Pt(1.5))

    add_tb(slide,
           left=Inches(0.5), top=Inches(1.38),
           width=Inches(5.55), height=Inches(0.42),
           text="Current: Synthetic Data Generation",
           size=Pt(13), bold=True, color=C_AMBER)

    synth_label = "SYNTHETIC DATA — not real sensor readings"
    add_shape(slide,
              left=Inches(0.5), top=Inches(1.83),
              width=Inches(5.55), height=Inches(0.28),
              fill=RGBColor(0x3A, 0x1A, 0x00))
    add_tb(slide,
           left=Inches(0.52), top=Inches(1.84),
           width=Inches(5.5), height=Inches(0.26),
           text=synth_label, size=Pt(9.5), bold=True,
           color=C_AMBER, align=PP_ALIGN.CENTER)

    synth_lines = [
        "Why synthetic data was required:",
        "",
        "  Real sonar hardware and calibrated marine sensor",
        "  arrays have not yet been deployed for this project.",
        "",
        "  Expert-labelled field surveys take months to",
        "  collect across four remote Indian reef regions.",
        "",
        "  Synthetic data lets the full MLOps pipeline be",
        "  built, tested, and validated before physical",
        "  hardware is available.",
        "",
        "How it was generated:",
        "  15,000 observations generated by src/data/generate_data.py",
        "  Statistical rules from published Indian reef literature",
        "  (Hughes et al. 2017; Venkataraman et al. 2011)",
        "  Gaussian noise (scale=0.15) prevents perfect separability",
        "  Fixed random seed 42 — fully reproducible",
    ]
    y = 2.18
    for line in synth_lines:
        bold = line.startswith("Why ") or line.startswith("How ")
        col  = C_CYAN if bold else C_TEXT
        add_tb(slide,
               left=Inches(0.48), top=Inches(y),
               width=Inches(5.6), height=Inches(0.27),
               text=line, size=Pt(10), color=col, bold=bold)
        y += 0.265

    # Right: proposed future
    add_shape(slide,
              left=Inches(6.55), top=Inches(1.3),
              width=Inches(6.45), height=Inches(5.6),
              fill=C_NAVY3, line=C_TEAL, line_w=Pt(1.0))

    add_tb(slide,
           left=Inches(6.75), top=Inches(1.38),
           width=Inches(6.1), height=Inches(0.42),
           text="Proposed Future: Real Field Collection",
           size=Pt(13), bold=True, color=C_TEAL)

    future_groups = [
        ("Acoustic Surveys", [
            "Multi-beam sonar deployed from survey vessels",
            "Measures bathymetric rugosity and backscatter",
            "Autonomous Underwater Vehicles (AUVs) for",
            "deep transects and consistent coverage",
        ], C_CYAN),
        ("Environmental Sensors", [
            "Fixed buoy arrays for temperature, pH,",
            "dissolved oxygen, turbidity, salinity",
            "Remotely Operated Vehicles (ROVs) for",
            "depth-stratified profiles",
        ], C_TEAL),
        ("Biological Surveys", [
            "Expert diver transects for coral cover,",
            "bleaching extent, and disease prevalence",
            "Photographic and video records for",
            "systematic label assignment",
        ], C_GREEN),
    ]

    y = 1.85
    for group_title, group_lines, color in future_groups:
        add_tb(slide,
               left=Inches(6.73), top=Inches(y),
               width=Inches(6.1), height=Inches(0.3),
               text=group_title, size=Pt(11), bold=True, color=color)
        y += 0.32
        for line in group_lines:
            add_tb(slide,
                   left=Inches(6.9), top=Inches(y),
                   width=Inches(5.9), height=Inches(0.26),
                   text=line, size=Pt(10), color=C_TEXT)
            y += 0.255
        y += 0.12

    add_notes(slide,
        "No real sonar data or field surveys have been collected for this prototype. "
        "Real sonar hardware, environmental buoys, and biological survey teams have not "
        "yet been deployed. All 15,000 observations were generated using a documented "
        "Python script that applies statistical rules derived from published Indian reef "
        "literature. Gaussian noise prevents any single feature from being a perfect "
        "predictor. The key advantage is that the complete MLOps pipeline — validation, "
        "preprocessing, training, evaluation, serving — can be built and tested now, and "
        "then the synthetic generator simply gets replaced with a real ingestion module "
        "when field data becomes available. The proposed future collection plan involves "
        "multi-beam sonar from boats and AUVs, fixed environmental sensor buoys, and "
        "expert biological surveys conducted by trained marine biologists."
    )
    return slide


def slide_04_features(prs):
    """Data sources and feature groups."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Data Sources and Feature Groups",
                    "21 columns across five categories — synthetic prototype")
    add_footer(slide, 4)

    groups = [
        {
            "title": "Geographic & Observation Metadata",
            "color": C_MUTED,
            "features": ["timestamp", "latitude", "longitude", "region"],
            "note": "4 reef regions in Indian waters",
        },
        {
            "title": "Sonar-Derived Structural Features",
            "color": C_CYAN,
            "features": ["sonar_backscatter (dB)", "rugosity_index",
                         "hard_substrate_percentage (%)", "acoustic_complexity_index"],
            "note": "Sonar measures structure — NOT bleaching or coral cover",
        },
        {
            "title": "Environmental Sensor Features",
            "color": C_TEAL,
            "features": ["water_temperature_c", "ph", "salinity_ppt",
                         "dissolved_oxygen_mg_l", "turbidity_ntu",
                         "light_intensity (PAR)", "current_speed_m_s", "depth_m"],
            "note": "Physical oceanographic conditions",
        },
        {
            "title": "Reef-Condition Observations",
            "color": C_GREEN,
            "features": ["coral_cover_percentage (%)",
                         "bleaching_percentage (%)",
                         "disease_percentage (%)"],
            "note": "Biological survey observations",
        },
        {
            "title": "Target Labels",
            "color": C_CORAL,
            "features": ["reef_health (4 classes)",
                         "restoration_suitability (3 classes)"],
            "note": "Supervised learning targets",
        },
    ]

    positions = [
        (0.28, 1.25, 2.45, 5.62),   # geo
        (2.88, 1.25, 2.6,  5.62),   # sonar
        (5.62, 1.25, 2.7,  5.62),   # env sensors
        (8.45, 1.25, 2.65, 5.62),   # reef condition
        (11.22, 1.25, 1.88, 5.62),  # targets
    ]

    for grp, (lx, ty, bw, bh) in zip(groups, positions):
        # Card background
        add_shape(slide,
                  left=Inches(lx), top=Inches(ty),
                  width=Inches(bw), height=Inches(bh),
                  fill=C_NAVY3, line=grp["color"], line_w=Pt(1.2))

        # Color header strip
        add_shape(slide,
                  left=Inches(lx), top=Inches(ty),
                  width=Inches(bw), height=Inches(0.42),
                  fill=grp["color"])

        add_tb(slide,
               left=Inches(lx + 0.06), top=Inches(ty + 0.02),
               width=Inches(bw - 0.08), height=Inches(0.38),
               text=grp["title"],
               size=Pt(9.5), bold=True,
               color=C_NAVY if grp["color"] != C_MUTED else C_NAVY,
               align=PP_ALIGN.CENTER)

        y = ty + 0.52
        for feat in grp["features"]:
            add_tb(slide,
                   left=Inches(lx + 0.1), top=Inches(y),
                   width=Inches(bw - 0.16), height=Inches(0.32),
                   text=feat,
                   size=Pt(9.5), color=C_TEXT)
            y += 0.33

        # Note at bottom
        add_shape(slide,
                  left=Inches(lx + 0.06), top=Inches(ty + bh - 0.52),
                  width=Inches(bw - 0.12), height=Inches(0.44),
                  fill=C_HIGHLIGHT)
        add_tb(slide,
               left=Inches(lx + 0.08), top=Inches(ty + bh - 0.50),
               width=Inches(bw - 0.14), height=Inches(0.42),
               text=grp["note"],
               size=Pt(8.2), italic=True,
               color=grp["color"])

    # Scientific limitation note at bottom
    add_shape(slide,
              left=Inches(0.28), top=Inches(6.95),
              width=Inches(12.82), height=Inches(0.35),
              fill=RGBColor(0x2A, 0x0A, 0x08))
    add_tb(slide,
           left=Inches(0.35), top=Inches(6.96),
           width=Inches(12.7), height=Inches(0.32),
           text="Scientific limit: Sonar measures habitat structure (depth, backscatter, rugosity). "
                "It does NOT directly measure coral bleaching, coral cover, or water chemistry.",
           size=Pt(9.5), italic=True, color=C_CORAL)

    add_notes(slide,
        "The dataset has 21 columns divided into five logical groups. Geographic metadata "
        "locates each observation in one of four Indian reef regions. Sonar-derived features "
        "describe the physical structure of the reef — rugosity, backscatter, hard substrate "
        "fraction, and acoustic complexity. These are habitat structure indicators only. "
        "An important scientific limit is that sonar cannot directly measure bleaching "
        "percentages, coral cover, or water chemistry. Environmental sensors capture the "
        "physical conditions of the water column — temperature, pH, dissolved oxygen, "
        "turbidity, salinity, light, and current. Reef-condition observations record what a "
        "diver or camera would see — coral cover, bleaching extent, and disease prevalence. "
        "Finally, we have two target label columns. In a real deployment, the sonar and "
        "sensor data would be collected continuously, while biological survey labels would "
        "require periodic expert annotation."
    )
    return slide


def slide_05_dataset(prs):
    """Dataset description with compact visualization."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Dataset Description",
                    "15,000 synthetic observations — confirmed from data/raw/observations.csv")
    add_footer(slide, 5)

    # Stats grid (left panel)
    add_shape(slide,
              left=Inches(0.3), top=Inches(1.28),
              width=Inches(4.2), height=Inches(5.65),
              fill=C_NAVY3, line=C_TEAL, line_w=Pt(0.8))

    add_tb(slide,
           left=Inches(0.5), top=Inches(1.36),
           width=Inches(3.8), height=Inches(0.38),
           text="Dataset Statistics",
           size=Pt(13), bold=True, color=C_CYAN)

    stats = [
        ("Total observations",   "15,000"),
        ("CSV columns",          "21"),
        ("Training rows (80%)",  "12,000"),
        ("Test rows (20%)",      "3,000"),
        ("Missing values",       "0  (all required)"),
        ("Duplicate rows",       "None in generated data"),
        ("Temporal range",       "2018-01-01 to 2024-12-31"),
        ("Random seed",          "42  (reproducible)"),
        ("Noise scale",          "0.15  (prevents perfect accuracy)"),
        ("CV folds",             "5  (stratified)"),
        ("Reef regions",         "4  (Indian waters)"),
        ("Health classes",       "4"),
        ("Restoration classes",  "3"),
        ("Numeric features",     "15"),
        ("Categorical features", "1  (region)"),
        ("Derived features",     "6  (engineered)"),
    ]

    y = 1.83
    for label, val in stats:
        add_tb(slide,
               left=Inches(0.45), top=Inches(y),
               width=Inches(2.0), height=Inches(0.27),
               text=label, size=Pt(9.5), color=C_MUTED)
        add_tb(slide,
               left=Inches(2.5), top=Inches(y),
               width=Inches(1.9), height=Inches(0.27),
               text=val, size=Pt(9.5), bold=True, color=C_TEXT)
        y += 0.29

    # Class lists (middle panel)
    add_shape(slide,
              left=Inches(4.7), top=Inches(1.28),
              width=Inches(4.0), height=Inches(2.7),
              fill=C_NAVY3, line=C_CYAN, line_w=Pt(0.8))

    add_tb(slide,
           left=Inches(4.9), top=Inches(1.36),
           width=Inches(3.7), height=Inches(0.38),
           text="Reef Health Classes (4)",
           size=Pt(12), bold=True, color=C_CYAN)

    health_classes = [
        ("healthy",           C_GREEN),
        ("stressed",          C_AMBER),
        ("bleached",          C_CORAL),
        ("severely_degraded", RGBColor(0xCC, 0x33, 0x33)),
    ]
    y = 1.83
    for cls, col in health_classes:
        add_shape(slide,
                  left=Inches(4.85), top=Inches(y + 0.04),
                  width=Inches(0.14), height=Inches(0.18),
                  fill=col)
        add_tb(slide,
               left=Inches(5.05), top=Inches(y),
               width=Inches(3.5), height=Inches(0.28),
               text=cls, size=Pt(11), color=C_TEXT)
        y += 0.38

    add_shape(slide,
              left=Inches(4.7), top=Inches(4.12),
              width=Inches(4.0), height=Inches(2.4),
              fill=C_NAVY3, line=C_TEAL, line_w=Pt(0.8))

    add_tb(slide,
           left=Inches(4.9), top=Inches(4.20),
           width=Inches(3.7), height=Inches(0.38),
           text="Restoration Suitability Classes (3)",
           size=Pt(12), bold=True, color=C_TEAL)

    rest_classes = [
        ("suitable",             C_GREEN),
        ("moderately_suitable",  C_AMBER),
        ("unsuitable",           C_CORAL),
    ]
    y = 4.65
    for cls, col in rest_classes:
        add_shape(slide,
                  left=Inches(4.85), top=Inches(y + 0.04),
                  width=Inches(0.14), height=Inches(0.18),
                  fill=col)
        add_tb(slide,
               left=Inches(5.05), top=Inches(y),
               width=Inches(3.5), height=Inches(0.28),
               text=cls, size=Pt(11), color=C_TEXT)
        y += 0.38

    # Chart (right panel)
    chart_bytes = make_dataset_summary_chart()
    add_image_bytes(slide, chart_bytes,
                    left=Inches(8.85), top=Inches(1.28),
                    width=Inches(4.2), height=Inches(5.65))

    add_notes(slide,
        "The dataset has exactly 15,000 rows and 21 columns, confirmed from the CSV file. "
        "The 80/20 stratified split gives 12,000 training rows and 3,000 test rows. "
        "Stratification ensures that class proportions are preserved in both splits. "
        "There are no missing values because all fields are required in the Pandera schema — "
        "any row with a null value would be rejected at validation. The class distribution is "
        "not perfectly balanced. For reef health, the stressed class is by far the most common, "
        "followed by severely degraded, healthy, and bleached. For restoration suitability, "
        "moderately suitable is the majority class. This imbalance is handled during training "
        "using the class_weight=balanced setting, which automatically adjusts the loss function "
        "to give more weight to minority classes. The charts shown here are approximate "
        "distributions scaled from the test-set support counts in the evaluation JSON files."
    )
    return slide


def slide_06_validation(prs):
    """Data quality and validation."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Data Quality and Validation",
                    "Pandera schema validation — pipeline gate before any modelling begins")
    add_footer(slide, 6)

    # Validation flow image
    val_img = make_validation_flow()
    add_image_bytes(slide, val_img,
                    left=Inches(0.25), top=Inches(1.22),
                    width=Inches(12.85), height=Inches(3.5))

    # Checks table
    checks = [
        ("Required columns",    "All 21 columns must be present"),
        ("Data types",          "Pandera coerces and enforces float64, datetime, str"),
        ("Sensor bounds",       "e.g. temperature 10–42 °C, pH 7.0–9.0, depth 0–50 m"),
        ("Missing values",      "All fields nullable=False — any null row fails"),
        ("Duplicate rows",      "Cross-field check, prevents repeated identical obs."),
        ("Coordinate bounds",   "Latitude/longitude must fall within the region bounding box"),
        ("Target labels",       "reef_health must be one of 4 valid classes; restoration one of 3"),
        ("Pipeline action",     "Validation failure stops the DVC pipeline — no training occurs"),
    ]

    col_widths = [2.5, 9.2]
    y = 4.88
    for i, (check, desc) in enumerate(checks):
        bg = C_NAVY3 if i % 2 == 0 else C_NAVY2
        add_shape(slide,
                  left=Inches(0.28), top=Inches(y),
                  width=Inches(12.75), height=Inches(0.295),
                  fill=bg)
        add_tb(slide,
               left=Inches(0.38), top=Inches(y + 0.02),
               width=Inches(col_widths[0]), height=Inches(0.26),
               text=check, size=Pt(10), bold=True, color=C_CYAN)
        add_tb(slide,
               left=Inches(col_widths[0] + 0.55), top=Inches(y + 0.02),
               width=Inches(col_widths[1]), height=Inches(0.26),
               text=desc, size=Pt(10), color=C_TEXT)
        y += 0.30

    add_notes(slide,
        "Data validation is the first gate in the DVC pipeline. Before any modelling "
        "work begins, every observation must pass a Pandera DataFrameModel schema check. "
        "Pandera is a Python library that lets you declare constraints on a DataFrame as "
        "a class definition, similar to Pydantic for API models. In lazy mode, it collects "
        "all validation errors across all rows before reporting them — so you get a complete "
        "picture of what went wrong rather than stopping at the first failure. "
        "The key checks are: all 21 columns must be present with the correct types, "
        "every sensor reading must fall within physically plausible bounds, no missing values "
        "are allowed, and target labels must match the known class sets. There is also a "
        "cross-field check that verifies each observation's coordinates fall within the "
        "bounding box of its declared region. If any of these checks fail, the pipeline "
        "exits with an error code and no training runs."
    )
    return slide


def slide_07_preprocessing(prs):
    """Preprocessing and feature transformation."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Preprocessing and Feature Transformation",
                    "Training-only preprocessor fitting — preventing test-data leakage")
    add_footer(slide, 7)

    # Flow diagram image
    flow_img = make_preprocessing_diagram()
    add_image_bytes(slide, flow_img,
                    left=Inches(0.25), top=Inches(1.22),
                    width=Inches(12.85), height=Inches(3.8))

    # Two detail panels below
    left_items = [
        ("Numerical Pipeline (15 features)", C_CYAN, [
            "1. SimpleImputer — median strategy (robust to skew)",
            "2. StandardScaler — zero mean, unit variance",
            "Handles: depth, temperature, pH, salinity, DO,",
            "turbidity, light, current, sonar, rugosity,",
            "hard_substrate, ACI, coral cover, bleaching, disease",
        ]),
    ]
    right_items = [
        ("Categorical Pipeline (1 feature: region)", C_TEAL, [
            "1. SimpleImputer — most_frequent strategy",
            "2. OneHotEncoder — 4 region values",
            "   Lakshadweep, Gulf of Mannar,",
            "   Gulf of Kutch, Andaman & Nicobar Islands",
            "handle_unknown='ignore' for inference safety",
        ]),
    ]

    for (title, color, lines), x_start in [
        (left_items[0], 0.28), (right_items[0], 6.65)
    ]:
        add_shape(slide,
                  left=Inches(x_start), top=Inches(5.15),
                  width=Inches(6.1), height=Inches(1.68),
                  fill=C_NAVY3, line=color, line_w=Pt(0.8))
        add_tb(slide,
               left=Inches(x_start + 0.15), top=Inches(5.22),
               width=Inches(5.8), height=Inches(0.32),
               text=title, size=Pt(11), bold=True, color=color)
        y = 5.58
        for line in lines:
            add_tb(slide,
                   left=Inches(x_start + 0.15), top=Inches(y),
                   width=Inches(5.8), height=Inches(0.26),
                   text=line, size=Pt(9.5), color=C_TEXT)
            y += 0.255

    # Leakage note
    add_shape(slide,
              left=Inches(0.28), top=Inches(6.91),
              width=Inches(12.75), height=Inches(0.35),
              fill=C_HIGHLIGHT)
    add_tb(slide,
           left=Inches(0.38), top=Inches(6.93),
           width=Inches(12.55), height=Inches(0.3),
           text="Preprocessor fitted on training split ONLY — test data statistics never influence scaling parameters.",
           size=Pt(10), bold=True, color=C_CORAL)

    add_notes(slide,
        "The preprocessing step converts raw feature columns into a numeric matrix that "
        "sklearn classifiers can consume. There are two separate pipelines: one for numeric "
        "features and one for the categorical region column. The critical design rule is "
        "that the ColumnTransformer is fitted only on the training split. If we fitted on "
        "the full dataset, the scaler would learn the mean and variance of test rows, "
        "which would constitute data leakage and inflate apparent test performance. "
        "Each task — health and restoration — gets its own fitted preprocessor saved as "
        "a .joblib file. These artifacts are used both during the DVC pipeline evaluation "
        "and at inference time in the FastAPI service. The preprocessor is always called "
        "in transform-only mode at inference time — never refit on production data. "
        "In addition to the 16 raw features, 6 derived features are computed before "
        "splitting: thermal stress index, oxygen stress index, acidity deviation, "
        "water quality index, substrate stability score, and structural complexity score."
    )
    return slide


def slide_08_training_design(prs):
    """Model training design."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Model Training Design",
                    "Three algorithms, two tasks, five-fold cross-validation — tracked in MLflow")
    add_footer(slide, 8)

    # Algorithm cards
    algorithms = [
        {
            "name": "Logistic Regression",
            "color": C_CYAN,
            "params": [
                "C = 1.0  (regularisation strength)",
                "max_iter = 1000",
                "class_weight = balanced",
                "solver = lbfgs  (default)",
            ],
            "why": "Strong baseline for linearly separable features. Interpretable coefficients.",
        },
        {
            "name": "Random Forest",
            "color": C_TEAL,
            "params": [
                "n_estimators = 200",
                "max_depth = 12",
                "min_samples_leaf = 4",
                "class_weight = balanced",
            ],
            "why": "Handles non-linear interactions between sensor readings. Robust to noise.",
        },
        {
            "name": "XGBoost",
            "color": C_AMBER,
            "params": [
                "n_estimators = 300",
                "max_depth = 6",
                "learning_rate = 0.05",
                "subsample = 0.8  |  colsample_bytree = 0.8",
            ],
            "why": "Gradient boosting often outperforms RF on tabular data with mixed features.",
        },
    ]

    for i, alg in enumerate(algorithms):
        lx = 0.3 + i * 4.28
        add_shape(slide,
                  left=Inches(lx), top=Inches(1.28),
                  width=Inches(4.0), height=Inches(4.0),
                  fill=C_NAVY3, line=alg["color"], line_w=Pt(1.5))
        # Header
        add_shape(slide,
                  left=Inches(lx), top=Inches(1.28),
                  width=Inches(4.0), height=Inches(0.42),
                  fill=alg["color"])
        add_tb(slide,
               left=Inches(lx + 0.1), top=Inches(1.30),
               width=Inches(3.8), height=Inches(0.38),
               text=alg["name"],
               size=Pt(12.5), bold=True,
               color=C_NAVY, align=PP_ALIGN.CENTER)
        add_tb(slide,
               left=Inches(lx + 0.1), top=Inches(1.76),
               width=Inches(3.8), height=Inches(0.28),
               text="Hyperparameters (params.yaml):",
               size=Pt(9.5), bold=True, color=C_MUTED)
        y = 2.08
        for param in alg["params"]:
            add_tb(slide,
                   left=Inches(lx + 0.18), top=Inches(y),
                   width=Inches(3.72), height=Inches(0.27),
                   text=param, size=Pt(10), color=C_TEXT)
            y += 0.275

        add_shape(slide,
                  left=Inches(lx + 0.1), top=Inches(y + 0.05),
                  width=Inches(3.8), height=Inches(0.025),
                  fill=alg["color"])
        y += 0.12
        add_tb(slide,
               left=Inches(lx + 0.1), top=Inches(y),
               width=Inches(3.8), height=Inches(0.68),
               text=alg["why"], size=Pt(9.5), italic=True, color=C_MUTED)

    # Training design details below
    design_items = [
        ("Separate training for each task",
         "Health and restoration are independent classifiers — not multi-output. Each trains on its own label column."),
        ("5-fold stratified cross-validation",
         "CV macro-F1 is the primary selection metric. Winner is selected before evaluation on the held-out test set."),
        ("class_weight = balanced",
         "Automatically adjusts loss function weights to compensate for class imbalance. All three algorithms apply this."),
        ("Fixed random seed = 42",
         "All randomness (data splits, model initialisation, numpy) is seeded for full reproducibility."),
        ("MLflow experiment tracking",
         "Every run logs: parameters, CV scores, test metrics, confusion matrix, feature importances, model artifact."),
    ]

    y = 5.38
    for i, (label, detail) in enumerate(design_items):
        bg = C_NAVY3 if i % 2 == 0 else C_NAVY2
        add_shape(slide,
                  left=Inches(0.28), top=Inches(y),
                  width=Inches(12.75), height=Inches(0.30),
                  fill=bg)
        add_tb(slide,
               left=Inches(0.38), top=Inches(y + 0.02),
               width=Inches(3.1), height=Inches(0.27),
               text=label, size=Pt(10), bold=True, color=C_CYAN)
        add_tb(slide,
               left=Inches(3.55), top=Inches(y + 0.02),
               width=Inches(9.4), height=Inches(0.27),
               text=detail, size=Pt(10), color=C_TEXT)
        y += 0.31

    add_notes(slide,
        "Three algorithms are trained for each classification task: Logistic Regression, "
        "Random Forest, and XGBoost. All hyperparameters are stored in params.yaml so they "
        "can be tracked and changed reproducibly via DVC. The two tasks are trained completely "
        "independently — health and restoration have separate preprocessors, separate training "
        "runs, and separate registered models. Class imbalance is handled by setting "
        "class_weight to balanced in all three algorithms. The primary selection metric is "
        "CV macro-F1, which averages F1 score equally across all classes regardless of their "
        "frequency. This matters because naive accuracy would favor the majority class. "
        "Five-fold cross-validation is used to estimate generalisation performance before "
        "any test data is touched. The model with the highest CV macro-F1 is selected, and "
        "only then evaluated on the held-out 3,000-row test set. MLflow tracks every experiment "
        "run — parameters, metrics, the trained model artifact, and confusion matrix — all "
        "stored in a local SQLite database at artifacts/mlruns.db."
    )
    return slide


def slide_09_health_results(prs):
    """Reef-health model comparison."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Reef-Health Model Comparison",
                    "Selected model: Logistic Regression (highest CV Macro-F1) — SYNTHETIC DATA")
    add_footer(slide, 9)

    chart_img = make_model_comparison_chart(
        HEALTH_METRICS,
        "Task: Reef Health Classification  |  4 classes  |  Test set = 3,000 rows"
    )
    add_image_bytes(slide, chart_img,
                    left=Inches(0.25), top=Inches(1.22),
                    width=Inches(9.35), height=Inches(4.5))

    C_GOLD = RGBColor(0xFF, 0xD7, 0x00)

    # Right panel — selected model details
    add_shape(slide,
              left=Inches(9.72), top=Inches(1.22),
              width=Inches(3.38), height=Inches(4.5),
              fill=C_NAVY3, line=C_GOLD, line_w=Pt(1.8))

    add_shape(slide,
              left=Inches(9.72), top=Inches(1.22),
              width=Inches(3.38), height=Inches(0.42),
              fill=RGBColor(0xFF, 0xD7, 0x00))
    add_tb(slide,
           left=Inches(9.8), top=Inches(1.24),
           width=Inches(3.2), height=Inches(0.38),
           text="Selected: Logistic Regression",
           size=Pt(10.5), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    selected_stats = [
        ("CV Macro-F1",         "0.7612  ± 0.009"),
        ("Test Macro-F1",       "0.7871"),
        ("Balanced Accuracy",   "0.8012"),
        ("Overall Accuracy",    "0.8193"),
        ("",                    ""),
        ("Per-class Test F1:",  ""),
        ("  healthy",           "0.6962"),
        ("  stressed",          "0.8986"),
        ("  bleached",          "0.8830"),
        ("  severely_degraded", "0.6708"),
    ]
    y = 1.73
    for label, val in selected_stats:
        lc = C_MUTED if label.startswith("  ") else C_TEXT
        vc = C_CYAN if val and not label.startswith("Per") else C_TEXT
        add_tb(slide,
               left=Inches(9.82), top=Inches(y),
               width=Inches(2.0), height=Inches(0.26),
               text=label, size=Pt(9.5), bold=(not label.startswith(" ")), color=lc)
        if val:
            add_tb(slide,
                   left=Inches(11.85), top=Inches(y),
                   width=Inches(1.15), height=Inches(0.26),
                   text=val, size=Pt(9.5), bold=True, color=vc, align=PP_ALIGN.RIGHT)
        y += 0.29

    # Per-class detail below chart
    add_shape(slide,
              left=Inches(0.25), top=Inches(5.82),
              width=Inches(12.85), height=Inches(1.02),
              fill=C_NAVY3, line=C_TEAL, line_w=Pt(0.5))

    add_tb(slide,
           left=Inches(0.4), top=Inches(5.88),
           width=Inches(12.5), height=Inches(0.28),
           text="Why Logistic Regression won: Although all three models achieved very similar CV macro-F1 "
                "(within 0.003 of each other), Logistic Regression scored highest at 0.7612 and also "
                "achieved the best test macro-F1 (0.7871) and balanced accuracy (0.8012).",
           size=Pt(10), color=C_TEXT)

    add_shape(slide,
              left=Inches(0.25), top=Inches(6.9),
              width=Inches(12.85), height=Inches(0.3),
              fill=C_HIGHLIGHT)
    add_tb(slide,
           left=Inches(0.35), top=Inches(6.92),
           width=Inches(12.65), height=Inches(0.26),
           text="SYNTHETIC DATA DISCLAIMER: These metrics reflect performance on synthetic observations only. "
                "They do not indicate real-world coral reef prediction accuracy.",
           size=Pt(9.5), italic=True, color=C_AMBER)

    add_notes(slide,
        "For the reef health task, all three algorithms performed very similarly. Logistic "
        "Regression achieved the highest five-fold cross-validation macro-F1 of 0.7612, "
        "slightly ahead of Random Forest at 0.7605 and XGBoost at 0.7580. On the held-out "
        "test set, Logistic Regression also achieved the best macro-F1 of 0.7871 and "
        "balanced accuracy of 0.8012. The stressed class is the easiest to classify, with "
        "an F1 of about 0.90 across all models — this is the majority class and has the "
        "clearest feature signal. The healthy and severely degraded classes are harder, "
        "both scoring around 0.67 to 0.70 F1. This reflects overlapping feature ranges "
        "between adjacent ordinal health states, which is intentional in the synthetic "
        "data generator. It is critical to remember that these scores are on synthetic "
        "data with known statistical structure. Real reef data would likely produce "
        "different results due to genuine environmental complexity and labelling noise."
    )
    return slide


def slide_10_restoration_results(prs):
    """Restoration model comparison."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Restoration Suitability Model Comparison",
                    "Selected model: XGBoost (highest CV Macro-F1) — SYNTHETIC DATA")
    add_footer(slide, 10)

    chart_img = make_model_comparison_chart(
        RESTORATION_METRICS,
        "Task: Restoration Suitability  |  3 classes  |  Test set = 3,000 rows"
    )
    add_image_bytes(slide, chart_img,
                    left=Inches(0.25), top=Inches(1.22),
                    width=Inches(9.35), height=Inches(4.5))

    # Right panel
    C_GOLD = RGBColor(0xFF, 0xD7, 0x00)
    add_shape(slide,
              left=Inches(9.72), top=Inches(1.22),
              width=Inches(3.38), height=Inches(4.5),
              fill=C_NAVY3, line=C_GOLD, line_w=Pt(1.8))
    add_shape(slide,
              left=Inches(9.72), top=Inches(1.22),
              width=Inches(3.38), height=Inches(0.42),
              fill=C_GOLD)
    add_tb(slide,
           left=Inches(9.8), top=Inches(1.24),
           width=Inches(3.2), height=Inches(0.38),
           text="Selected: XGBoost",
           size=Pt(10.5), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    selected_stats = [
        ("CV Macro-F1",         "0.7913  ± 0.008"),
        ("Test Macro-F1",       "0.8029"),
        ("Balanced Accuracy",   "0.8121"),
        ("Overall Accuracy",    "0.8277"),
        ("",                    ""),
        ("Per-class Test F1:",  ""),
        ("  suitable",          "0.7204"),
        ("  mod. suitable",     "0.8884"),
        ("  unsuitable",        "0.8000"),
    ]
    y = 1.73
    for label, val in selected_stats:
        lc = C_MUTED if label.startswith("  ") else C_TEXT
        vc = C_CYAN if val and not label.startswith("Per") else C_TEXT
        add_tb(slide,
               left=Inches(9.82), top=Inches(y),
               width=Inches(2.0), height=Inches(0.27),
               text=label, size=Pt(9.5), bold=(not label.startswith(" ")), color=lc)
        if val:
            add_tb(slide,
                   left=Inches(11.85), top=Inches(y),
                   width=Inches(1.15), height=Inches(0.27),
                   text=val, size=Pt(9.5), bold=True, color=vc, align=PP_ALIGN.RIGHT)
        y += 0.31

    add_shape(slide,
              left=Inches(0.25), top=Inches(5.82),
              width=Inches(12.85), height=Inches(1.02),
              fill=C_NAVY3, line=C_TEAL, line_w=Pt(0.5))

    add_tb(slide,
           left=Inches(0.4), top=Inches(5.88),
           width=Inches(12.5), height=Inches(0.28),
           text="Why XGBoost won: XGBoost achieved CV macro-F1 of 0.7913, ahead of Random Forest (0.7839) "
                "and Logistic Regression (0.7808). It also produced the best test macro-F1 (0.8029) "
                "and accuracy (0.8277), capturing non-linear interactions between substrate and water quality.",
           size=Pt(10), color=C_TEXT)

    add_shape(slide,
              left=Inches(0.25), top=Inches(6.9),
              width=Inches(12.85), height=Inches(0.3),
              fill=C_HIGHLIGHT)
    add_tb(slide,
           left=Inches(0.35), top=Inches(6.92),
           width=Inches(12.65), height=Inches(0.26),
           text="SYNTHETIC DATA DISCLAIMER: These metrics reflect performance on synthetic observations only. "
                "They do not indicate real-world coral reef prediction accuracy.",
           size=Pt(9.5), italic=True, color=C_AMBER)

    add_notes(slide,
        "For the restoration suitability task, XGBoost achieved the highest cross-validation "
        "macro-F1 of 0.7913, compared to 0.7839 for Random Forest and 0.7808 for Logistic "
        "Regression. On the test set, XGBoost also led with a macro-F1 of 0.8029 and "
        "accuracy of 0.8277. The moderately suitable class has the highest per-class F1 "
        "at 0.8884, which makes sense as it is the majority class. The suitable class is "
        "harder to classify at 0.7204, which reflects overlap with moderately suitable sites "
        "where most conditions are good but one or two factors are suboptimal. The unsuitable "
        "class achieves a solid 0.8000 F1, as the extreme conditions that make a site "
        "unsuitable are generally quite distinct. Comparing across the two tasks, XGBoost "
        "did better on restoration but not on health, suggesting the restoration decision "
        "depends more heavily on non-linear interactions — particularly between substrate "
        "quality, turbidity, and temperature — that boosted trees handle well."
    )
    return slide


def slide_11_reproducibility(prs):
    """Reproducibility and experiment tracking."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Reproducibility and Experiment Tracking",
                    "Every step from data generation to candidate registration is versioned")
    add_footer(slide, 11)

    # Pipeline diagram
    pipe_img = make_pipeline_diagram()
    add_image_bytes(slide, pipe_img,
                    left=Inches(0.25), top=Inches(1.22),
                    width=Inches(12.85), height=Inches(2.45))

    # Two column table of tools
    tools = [
        ("Git",             "Source code versioning — every training run is tied to a specific commit hash."),
        ("DVC",             "Data and pipeline versioning — dvc repro reruns only changed stages (params.yaml-driven)."),
        ("params.yaml",     "Single source of truth for all hyperparameters, dataset size, split ratios, thresholds."),
        ("MLflow",          "Experiment tracking — logs parameters, CV scores, test metrics, confusion matrix, model artifact."),
        ("Saved preprocessors",
                            "Fitted ColumnTransformer (.joblib) artifacts preserved — inference uses same scaling as training."),
        ("Model registry",  "Champions registered in MLflow registry with 'champion' alias. Candidates never auto-promoted."),
        ("1,023 tests",     "Full pytest suite (tests/) validates pipeline, API, drift, retraining, and governance."),
    ]

    y = 3.82
    for i, (tool, desc) in enumerate(tools):
        bg = C_NAVY3 if i % 2 == 0 else C_NAVY2
        add_shape(slide,
                  left=Inches(0.28), top=Inches(y),
                  width=Inches(12.75), height=Inches(0.415),
                  fill=bg)
        add_tb(slide,
               left=Inches(0.38), top=Inches(y + 0.06),
               width=Inches(2.3), height=Inches(0.3),
               text=tool, size=Pt(10.5), bold=True, color=C_CYAN)
        add_tb(slide,
               left=Inches(2.75), top=Inches(y + 0.06),
               width=Inches(10.22), height=Inches(0.3),
               text=desc, size=Pt(10.5), color=C_TEXT)
        y += 0.425

    add_notes(slide,
        "One of the core requirements of a production MLOps pipeline is reproducibility. "
        "If you run dvc repro with the same params.yaml and the same input data, you should "
        "get exactly the same models every time. This is achieved through several layers. "
        "Git versions the source code. DVC versions the data files and pipeline stages — "
        "if params.yaml hasn't changed and the data hasn't changed, DVC skips rerunning "
        "the expensive stages. params.yaml is the single source of truth — changing "
        "n_samples or a hyperparameter triggers only the affected downstream stages. "
        "MLflow records every training run with its full parameter set and metrics, so "
        "you can always go back and see exactly what configuration produced a given model. "
        "Fitted preprocessors are saved separately so inference always uses the same scaling "
        "statistics as training. Critically, candidate models are registered in MLflow but "
        "the champion alias is never moved automatically — explicit human approval is required "
        "to promote a new champion. This is covered in the model governance milestone."
    )
    return slide


def slide_12_conclusions(prs):
    """Current results, limitations and next steps."""
    slide = blank_slide(prs)
    set_bg(slide, C_NAVY)
    add_slide_title(slide, "Results, Limitations and Next Steps",
                    "Prototype performance on synthetic data — real-world validation required")
    add_footer(slide, 12)

    # Results panel (left)
    add_shape(slide,
              left=Inches(0.3), top=Inches(1.28),
              width=Inches(4.2), height=Inches(5.6),
              fill=C_NAVY3, line=C_GREEN, line_w=Pt(1.0))
    add_shape(slide,
              left=Inches(0.3), top=Inches(1.28),
              width=Inches(4.2), height=Inches(0.42),
              fill=C_GREEN)
    add_tb(slide,
           left=Inches(0.4), top=Inches(1.30),
           width=Inches(4.0), height=Inches(0.38),
           text="Selected Models",
           size=Pt(12), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    results = [
        ("Reef Health", "Logistic Regression"),
        ("CV Macro-F1", "0.7612"),
        ("Test Macro-F1", "0.7871"),
        ("Test Accuracy", "0.8193"),
        ("", ""),
        ("Restoration", "XGBoost"),
        ("CV Macro-F1", "0.7913"),
        ("Test Macro-F1", "0.8029"),
        ("Test Accuracy", "0.8277"),
    ]
    y = 1.82
    for label, val in results:
        if not label and not val:
            y += 0.12
            continue
        bold = label in ("Reef Health", "Restoration")
        col  = C_CYAN if bold else C_MUTED
        add_tb(slide,
               left=Inches(0.42), top=Inches(y),
               width=Inches(2.0), height=Inches(0.28),
               text=label, size=Pt(10.5), bold=bold, color=col)
        add_tb(slide,
               left=Inches(2.45), top=Inches(y),
               width=Inches(1.95), height=Inches(0.28),
               text=val, size=Pt(10.5), bold=bold,
               color=C_TEXT if not bold else C_CYAN, align=PP_ALIGN.RIGHT)
        y += 0.30

    # Limitations panel (middle)
    add_shape(slide,
              left=Inches(4.75), top=Inches(1.28),
              width=Inches(4.2), height=Inches(5.6),
              fill=C_NAVY3, line=C_CORAL, line_w=Pt(1.0))
    add_shape(slide,
              left=Inches(4.75), top=Inches(1.28),
              width=Inches(4.2), height=Inches(0.42),
              fill=C_CORAL)
    add_tb(slide,
           left=Inches(4.85), top=Inches(1.30),
           width=Inches(4.0), height=Inches(0.38),
           text="Current Limitations",
           size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    limitations = [
        "All 15,000 observations are synthetic.",
        "Labels are generated by a known scoring",
        "function — real labels require expert",
        "biological field surveys.",
        "",
        "Sonar features in this dataset are",
        "statistically simulated, not calibrated",
        "from real acoustic hardware.",
        "",
        "Metrics demonstrate prototype pipeline",
        "correctness, not real-world reef",
        "prediction accuracy.",
        "",
        "No temporal or regional hold-out splits",
        "— spatial autocorrelation not yet modelled.",
        "",
        "No independent external validation set",
        "from a separate reef system.",
    ]
    y = 1.82
    for line in limitations:
        add_tb(slide,
               left=Inches(4.9), top=Inches(y),
               width=Inches(3.9), height=Inches(0.265),
               text=line, size=Pt(10), color=C_TEXT)
        y += 0.265

    # Next steps panel (right)
    add_shape(slide,
              left=Inches(9.2), top=Inches(1.28),
              width=Inches(3.88), height=Inches(5.6),
              fill=C_NAVY3, line=C_CYAN, line_w=Pt(1.0))
    add_shape(slide,
              left=Inches(9.2), top=Inches(1.28),
              width=Inches(3.88), height=Inches(0.42),
              fill=C_CYAN)
    add_tb(slide,
           left=Inches(9.3), top=Inches(1.30),
           width=Inches(3.68), height=Inches(0.38),
           text="Next Steps",
           size=Pt(12), bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    next_steps = [
        ("Field data integration", [
            "Deploy real sonar hardware and",
            "calibrated environmental sensors.",
            "Commission expert biological surveys",
            "for ground-truth label annotation.",
        ]),
        ("Validation strategy", [
            "Hold out a geographically separate",
            "reef region for external evaluation.",
            "Temporal cross-validation across",
            "survey years (2018–2024).",
        ]),
        ("Model improvement", [
            "SHAP feature importance analysis",
            "to identify key sensor drivers.",
            "Hyperparameter search beyond",
            "current params.yaml defaults.",
        ]),
    ]
    y = 1.82
    for step_title, step_lines in next_steps:
        add_tb(slide,
               left=Inches(9.3), top=Inches(y),
               width=Inches(3.68), height=Inches(0.28),
               text=step_title, size=Pt(10.5), bold=True, color=C_CYAN)
        y += 0.30
        for line in step_lines:
            add_tb(slide,
                   left=Inches(9.4), top=Inches(y),
                   width=Inches(3.58), height=Inches(0.255),
                   text=line, size=Pt(9.5), color=C_TEXT)
            y += 0.255
        y += 0.14

    # Conclusion bar
    add_shape(slide,
              left=Inches(0.3), top=Inches(6.96),
              width=Inches(12.78), height=Inches(0.32),
              fill=C_NAVY2)
    add_tb(slide,
           left=Inches(0.4), top=Inches(6.97),
           width=Inches(12.6), height=Inches(0.28),
           text="The pipeline is complete and test-verified. Replacing synthetic data with real calibrated "
                "field observations is the critical next milestone before any ecological conclusions can be drawn.",
           size=Pt(10), italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide,
        "To summarise the training results: for reef health classification, Logistic Regression "
        "was selected with a test macro-F1 of 0.7871 and accuracy of 0.8193. For restoration "
        "suitability, XGBoost was selected with a test macro-F1 of 0.8029 and accuracy of "
        "0.8277. Both models are registered in the MLflow model registry with the 'champion' "
        "alias. These are strong results for a classification prototype, but the critical "
        "limitation is that the data is entirely synthetic. The synthetic generator bakes in "
        "certain correlations by design, which means the models are learning the structure of "
        "the simulation rather than the complexity of real reef ecosystems. The next concrete "
        "steps are to deploy real hardware and collect labelled field data, and then to evaluate "
        "the models on a geographically separated hold-out reef to test for spatial generalisation. "
        "The complete MLOps infrastructure — the validation pipeline, drift monitoring, retraining "
        "governance, FastAPI inference service, and Streamlit dashboard — are already built and "
        "waiting for real data to flow through them."
    )
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    print("Creating presentation...")

    prs = new_prs()

    builders = [
        slide_01_title,
        slide_02_problem,
        slide_03_data_collection,
        slide_04_features,
        slide_05_dataset,
        slide_06_validation,
        slide_07_preprocessing,
        slide_08_training_design,
        slide_09_health_results,
        slide_10_restoration_results,
        slide_11_reproducibility,
        slide_12_conclusions,
    ]

    for i, builder in enumerate(builders, 1):
        print(f"  Building slide {i:02d}/{len(builders)} — {builder.__name__}...")
        builder(prs)

    prs.save(str(OUT_FILE))

    size_mb = OUT_FILE.stat().st_size / (1024 * 1024)
    print(f"\nSaved: {OUT_FILE}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size:   {size_mb:.2f} MB")

    # ── Quick structural verification ─────────────────────────────────────────
    print("\nVerifying structure...")
    prs2 = Presentation(str(OUT_FILE))
    assert len(prs2.slides) == 12, f"Expected 12 slides, got {len(prs2.slides)}"

    coralsense_hits = []
    for i, slide in enumerate(prs2.slides, 1):
        notes = slide.notes_slide.notes_text_frame.text
        assert len(notes) > 50, f"Slide {i} has insufficient speaker notes"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Check for user-facing "CoralSense" string
                        if "CoralSense" in run.text and i > 1:
                            coralsense_hits.append((i, run.text[:60]))

    if coralsense_hits:
        print(f"  WARNING: CoralSense found in slides: {coralsense_hits}")
    else:
        print("  No user-facing 'CoralSense' text found in slides 2-12")

    print(f"  All 12 slides verified")
    print(f"  Speaker notes present on all slides")
    print(f"  File size: {size_mb:.2f} MB {'(OK)' if size_mb < 10 else '(EXCEEDS 10 MB)'}")
    print("\nDone.")


if __name__ == "__main__":
    main()
